"""3차 업무 PPT 생성 — 코멘토 3차 업무 / 정승원

2차 업무 PPT 서식을 그대로 따릅니다.
  13.333 x 7.5 in · 배경 #FAFAFA · 카드 #FFFFFF + #EBEBEB 0.75pt 테두리
  제목 Pretendard SemiBold 27pt · 본문 Pretendard · 수치 Cascadia Mono

구성
----
  01 방법    핀란드 연안 데이터로 학습셋을 만들고 파인튜닝한 경로
  02 결과    대조군 넷의 기여도, 그리고 지표를 바꾸니 뒤집힌 결론
  03 확인    탐지 결과를 한 장으로 크게, 원본과 나란히

한국 확장은 아직 넣지 않습니다. 부산에서 잰 수치가 규약 차이에 크게 흔들려,
무엇을 어디까지 말할 수 있는지 정하기 전에는 장을 만들지 않습니다.
slide4 는 지울 때가 아니라 다시 쓸 때를 위해 남겨 둡니다.
"""
import os
import json
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

W, H = 13.333, 7.5
INK, BODY, MUTE, FAINT = "171717", "4D4D4D", "8F8F8F", "A1A1A1"
HAIR, CANVAS, CARD = "EBEBEB", "FAFAFA", "FFFFFF"
ACC, WARN = "0070F3", "D4443C"
SANS, SEMI, MED, MONO = ("Pretendard", "Pretendard SemiBold",
                         "Pretendard Medium", "Cascadia Mono")
TOTAL = 3


def C(h):
    return RGBColor.from_string(h)


def txt(sl, l, t, w, h, s, font=SANS, size=10, color=BODY,
        align=PP_ALIGN.LEFT, line=1.25):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = C(color)
    return tb


def card(sl, l, t, w, h, radius=0.045):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = C(CARD)
    sh.line.color.rgb = C(HAIR)
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if sh.has_text_frame:
        sh.text_frame.text = ""
    return sh


def base(prs, eyebrow, title, sub, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C(CANVAS)
    bg.line.fill.background()
    bg.shadow.inherit = False
    txt(sl, .72, .46, 11.89, .24, eyebrow, MONO, 8.5, MUTE)
    txt(sl, .72, .74, 11.89, .52, title, SEMI, 27, INK)
    txt(sl, .72, 1.28, 11.89, .26, sub, SANS, 10.5, BODY)
    txt(sl, 11.40, 6.92, 1.21, .24, "%02d / %02d" % (page, TOTAL),
        MONO, 8.5, FAINT, PP_ALIGN.RIGHT)
    return sl


def rows(sl, l, t, cols, data, widths, size=9.0):
    """텍스트박스 격자로 만든 간단한 표입니다. 셀이 (값, 색) 튜플이면 색을 적용합니다."""
    x = l
    for c, wd in zip(cols, widths):
        if c:
            txt(sl, x, t, wd, .2, c, MED, 8.5, MUTE)
        x += wd
    y = t + .26
    for r in data:
        x = l
        for ci, (v, wd) in enumerate(zip(r, widths)):
            col = BODY
            if isinstance(v, tuple):
                v, col = v
            txt(sl, x, y, wd, .2, str(v), SANS if ci == 0 else MONO, size, col)
            x += wd
        y += .245
    return y


def strip(sl, l, t, w, stats, gap=None):
    """카드 안에 큰 수치를 나란히 놓습니다."""
    gap = gap or w / len(stats)
    x = l
    for lab, val in stats:
        txt(sl, x, t, gap - .1, .22, lab, SANS, 9.0, MUTE)
        txt(sl, x, t + .30, gap - .1, .34, val, SEMI, 17, INK)
        x += gap


# ------------------------------------------------------------------ 슬라이드
def slide1(prs, fig):
    sl = base(prs, "01 / 방법",
              "5 픽셀짜리 선박 탐지 — 데이터와 학습",
              "Sentinel-2 GSD 10 m · 핀란드 연안 7,010 타일 / 13,069 인스턴스 · "
              "DOTA 사전학습 YOLO11s-OBB 파인튜닝 · 지역 단위 분할", 1)

    card(sl, .72, 1.70, 6.42, 2.94, .04)
    txt(sl, .98, 1.90, 5.90, .24, "주석만 공개된 데이터 → 학습셋",
        SEMI, 10.5, INK)
    txt(sl, .98, 2.22, 5.90, .22,
        "Zenodo 는 GeoPackage 주석만 줍니다. 사진은 따로 맞춰야 합니다.",
        SANS, 9.5, MUTE)
    txt(sl, .98, 2.60, 5.90, 1.70,
        "1.  Zenodo 15019034 주석 (CC-BY-4.0, RSE 2025)\n"
        "2.  AWS 공개 COG 에서 Sentinel-2 L2A TCI 수신 — 인증 불필요\n"
        "3.  주석 좌표를 COG 격자에 정합해 320 x 320 타일로 절단\n"
        "4.  DOTA 사전학습 YOLO11s-OBB 파인튜닝 (T4, 100 epoch)\n"
        "5.  학습에 안 쓴 해역(34WFT)에서만 평가",
        SANS, 9.5, BODY, line=1.52)
    txt(sl, .98, 4.28, 5.90, .24,
        "→ 주석은 L1C 기준이지만 L2A 는 같은 촬영·같은 격자라 기하가 같습니다",
        MED, 9.5, INK)

    card(sl, 7.36, 1.70, 5.25, 2.94, .04)
    txt(sl, 7.62, 1.90, 4.73, .24, "지역 단위 분할 — 수치를 믿을 근거",
        SEMI, 10.5, INK)
    txt(sl, 7.62, 2.22, 4.73, .22,
        "같은 해역이 학습과 평가에 섞이면 성능이 부풀려집니다", SANS, 9.5, MUTE)
    rows(sl, 7.62, 2.58, ["split", "타일", "인스턴스", "지역"],
         [["train", "6,261", "12,354", "34VEM 외"],
          ["val", "353", "349", "34VER"],
          ["test", "396", "366", "34WFT"]],
         [1.15, 1.02, 1.28, 1.28])
    txt(sl, 7.62, 3.52, 4.73, .84,
        "선박 크기는 20 m 부터 700 m 까지 넓습니다. 150 m 이상 대형선이 665척으로,\n"
        "\"레저용 소형선만 있을 것\" 이라는 예상은 측정으로 뒤집혔습니다.\n"
        "장변 중앙값은 52.3 m = 5.23 px 입니다.",
        SANS, 9.0, MUTE, line=1.42)

    card(sl, .72, 4.86, 11.89, 1.90, .08)
    txt(sl, .98, 5.08, 11.37, .24, "최종 성능 — 못 본 해역에서",
        SEMI, 10.5, INK)
    strip(sl, .98, 5.42, 11.37,
          [("test mAP50", "0.9014"), ("test mAP50-95", "0.5544"),
           ("Precision", "0.834"), ("Recall", "0.877"),
           ("CPU 타일 한 장", "54 ms")], 2.28)
    txt(sl, .98, 6.32, 11.37, .24,
        "Tesla T4 로 학습했고, 추론은 CPU 로도 됩니다 — 항만 한 곳 전체 주사가 1.61 초입니다",
        SANS, 9.5, MUTE)


def slide2(prs, fig):
    sl = base(prs, "02 / 결과",
              "성능의 출처와 지표의 함정",
              "대조군 넷으로 기여도를 갈랐습니다 · 같은 모델을 점 기준으로 재니 결론이 뒤집혔습니다", 2)

    card(sl, .72, 1.70, 6.42, 3.06, .04)
    txt(sl, .98, 1.90, 5.90, .24, "대조군 넷의 기여도",
        SEMI, 10.5, INK)
    txt(sl, .98, 2.22, 5.90, .22,
        "한 번에 한 축만 바꿔 기여도를 귀속시킵니다", SANS, 9.5, MUTE)
    rows(sl, .98, 2.58, ["실험", "val mAP50", "test mAP50", "Recall"],
         [["밑바닥 (사전학습 없음)", "0.8106", "0.8734", "0.831"],
          ["증강 off", ("0.7068", WARN), ("0.7678", WARN), ("0.735", WARN)],
          ["DOTA 사전학습 (n)", "0.8390", "0.8809", "0.828"],
          ["DOTA 사전학습 (s)", ("0.8628", ACC), ("0.9014", ACC), ("0.877", ACC)]],
         [2.30, 1.20, 1.24, 1.16])
    txt(sl, .98, 3.96, 5.90, .58,
        "증강  +13.2%p      사전학습  +2.8%p      모델 크기  +2.4%p",
        MONO, 10.0, INK, line=1.3)
    txt(sl, .98, 4.42, 5.90, .24,
        "→ 가장 크게 기여한 것은 구조가 아니라 증강이었습니다", MED, 9.5, INK)

    card(sl, 7.36, 1.70, 5.25, 3.06, .04)
    txt(sl, 7.62, 1.90, 4.73, .24, "지표에 따라 뒤집히는 결론",
        SEMI, 10.5, INK)
    txt(sl, 7.62, 2.22, 4.73, .22,
        "같은 모델, 같은 평가셋입니다. 재는 방법만 다릅니다", SANS, 9.5, MUTE)
    rows(sl, 7.62, 2.58, ["지표", "값"],
         [["IoU 기반 mAP50", "0.901"],
          ["IoU 기반 mAP50-95", ("0.459", WARN)],
          ["점 기반 AP", ("0.911", ACC)],
          ["점 기반 재현율", ("0.992", ACC)],
          ["중심 오차 (중앙값)", "0.46 px ≈ 5 m"]],
         [2.70, 2.03])
    txt(sl, 7.62, 4.06, 4.73, .58,
        "5 픽셀짜리 물체에 IoU 0.75 를 요구하는 것은\n"
        "부분 픽셀 정확도를 요구하는 것과 같습니다.",
        MED, 9.5, INK, line=1.4)

    card(sl, .72, 4.98, 11.89, 1.76, .085)
    txt(sl, .98, 5.18, 11.37, .24,
        "낮은 IoU 의 원인 — 라벨 좌표 전수 확인", SEMI, 10.5, INK)
    txt(sl, .98, 5.52, 11.37, 1.02,
        "여덟 좌표로 배포되기에 회전 박스라고 믿었는데, 13,069개 전부가 축 정렬이었습니다. "
        "장변 각도가 0~10도에 6,343개, 80~90도에 6,726개, 그 사이는 0개입니다.\n"
        "종횡비도 물리가 아니었습니다. 0.958 에서 0.493 으로 크기에 따라 단조 감소하는데, "
        "실제 선박 비율은 크기와 무관하므로 이것은 배의 모양이 아니라 작은 배를 뭉툭하게 칠한 주석 습관입니다.\n"
        "그리고 폭은 어떤 종횡비를 가정해도 72~95% 가 2 px 미만입니다 — 규약을 정하는 문제가 아니라 잴 대상이 없습니다.",
        SANS, 9.0, MUTE, line=1.42)


def slide3(prs, fig):
    # 수치는 그림을 만든 스크립트가 남긴 것을 그대로 읽습니다. 손으로 적으면
    # 그림을 다시 뽑았을 때 글과 그림이 어긋납니다.
    m = {"gt": 11, "det": 8, "hit": 6, "unmatched": 2, "conf": 0.5,
         "km": [6.4, 3.6], "dn": {"hit": 21, "miss": 20, "unmatched": 64}}
    jp = fig("fig6_hero.json")
    if os.path.exists(jp):
        m.update(json.load(open(jp, encoding="utf-8")))
    kw, kh = m["km"]
    dn = m["dn"]

    sl = base(prs, "03 / 확인", "탐지 결과 육안 확인",
              f"학습에 안 쓴 해역(34WFT)의 장면 원본 · 선박이 몰린 {kw} x {kh} km 를 "
              f"겹쳐 훑어 탐지 · 신뢰도 {m['conf']} (F1 최고점)", 3)

    # 이 장의 주장은 큰 그림 한 장입니다. 배지·범례·축척은 그림 안에
    # 이미 얹혀 있어, 슬라이드에서는 자리만 잡아 줍니다.
    card(sl, .64, 1.62, 8.71, 4.97, .02)
    p = fig("fig6_hero_det.png")
    if os.path.exists(p):
        sl.shapes.add_picture(p, Inches(.72), Inches(1.70), width=Inches(8.55))
    txt(sl, .72, 6.62, 8.55, .24,
        "배는 장변이 5 화소입니다. 형체가 읽히도록 6배로 늘렸고, 그래서 화소도 "
        "함께 보입니다 — 이것이 10 m 해상도에서 배가 보이는 실제 크기입니다.",
        MED, 9.5, INK)

    txt(sl, 9.57, 1.70, 2.96, .22, "같은 장면, 박스 없이", SEMI, 10.5, INK)
    txt(sl, 9.57, 1.96, 2.96, .20, "Sentinel-2 L2A TCI · 받은 그대로", SANS, 9.0, MUTE)
    card(sl, 9.49, 2.16, 3.12, 1.83, .03)
    p = fig("fig6_hero_raw.png")
    if os.path.exists(p):
        sl.shapes.add_picture(p, Inches(9.57), Inches(2.24), width=Inches(2.96))

    # 표는 빼기를 시키지 않습니다. "주석 11, 탐지 8" 만 놓으면 읽는 사람이
    # 11-8=3 을 미탐으로 셉니다. 탐지 8 중 2 는 주석에 없는 자리라 실제 미탐은
    # 5 입니다. 그래서 정답과 탐지를 각각 맞은 것/틀린 것으로 갈라 적습니다.
    txt(sl, 9.57, 4.20, 2.96, .22, "이 창의 숫자", SEMI, 10.5, INK)
    rows(sl, 9.57, 4.28, ["", ""],
         [["주석된 선박", f"{m['gt']}척"],
          ["  그중 찾은 것", f"{m['hit']}척"],
          ["  놓친 것", f"{m['gt'] - m['hit']}척"],
          ["주석에 없는데 잡은 것", f"{m['unmatched']}척"]],
         [1.70, 1.26])
    txt(sl, 9.57, 5.62, 2.96, 1.30,
        f"놓친 {m['gt'] - m['hit']}척은 주변 물과의 밝기 차가 {dn['miss']} DN 인 "
        f"어두운 배입니다. 주석에 없는데 잡은 {m['unmatched']}척은 {dn['unmatched']} DN "
        f"으로 찾은 배({dn['hit']} DN)보다 오히려 밝아, 주석이 빠뜨린 배로 보입니다.",
        SANS, 9.0, MUTE, line=1.42)


def slide4(prs, fig):
    sl = base(prs, "04 / 확장",
              "한국 항만으로 확장",
              "핀란드에서 검증한 모델을 한국 항만에 적용합니다 · 4차 웹앱의 밑그림", 4)

    card(sl, .72, 1.70, 6.42, 3.10, .04)
    txt(sl, .98, 1.90, 5.90, .24, "이미 준비된 것", SEMI, 10.5, INK)
    txt(sl, .98, 2.22, 5.90, .22,
        "항만 이름 → 무운 장면 검색 → 수신 → 겹쳐 탐지 → NMS", SANS, 9.5, MUTE)
    rows(sl, .98, 2.58, ["", "내용"],
         [["항만 목록", "22곳 등록 완료 — 먼저 한 곳부터"],
          ["장면 수신", "AWS 공개 COG · 인증 불필요"],
          ["추론", "CPU 만으로 항만 한 곳 1.61 초"],
          ["참조 자료", "GFW 한국 근해 443,344건"]],
         [1.90, 3.90])
    txt(sl, .98, 3.92, 5.90, .76,
        "GPU 서버를 상시 켜 둘 필요가 없습니다. 웹앱에서 항구를 고르면\n"
        "그 자리에서 최신 위성사진을 받아 탐지가 돌아갑니다.\n"
        "GFW 는 57.1% 가 AIS 로 대조돼 있어 독립적인 참조가 됩니다.",
        SANS, 9.0, MUTE, line=1.42)

    card(sl, 7.36, 1.70, 5.25, 3.10, .04)
    txt(sl, 7.62, 1.90, 4.73, .24, "먼저 한 곳부터", SEMI, 10.5, INK)
    txt(sl, 7.62, 2.22, 4.73, .22,
        "학습 조건과 닮은 항만부터 시작합니다", SANS, 9.5, MUTE)
    rows(sl, 7.62, 2.58, ["조건", "유리한 쪽"],
         [["배경", "열린 바다 · 부두와 갯벌이 적을 것"],
          ["물 색", "어두운 물"],
          ["선박", "정박 중인 대형선 · 형태 안정"],
          ["표본", "수십 척 이상"]],
         [1.10, 3.63])
    txt(sl, 7.62, 3.94, 4.73, .82,
        "핀란드 학습 데이터의 전제인 \"어두운 물 위의 밝은 덩어리\" 가\n"
        "그대로 성립하는 곳입니다.\n\n"
        "항구를 늘리는 일은 언제든 되지만, 한 곳을 제대로 만드는 것이\n"
        "먼저입니다.",
        SANS, 9.0, MUTE, line=1.42)

    card(sl, .72, 5.02, 11.89, 1.72, .085)
    txt(sl, .98, 5.22, 11.37, .24, "웹앱 구상과 그 다음", SEMI, 10.5, INK)
    txt(sl, .98, 5.56, 11.37, .82,
        "항구를 드롭다운에서 고르면 최신 무운 장면을 찾아 받고, 겹쳐 훑어 탐지한 뒤 지도에 표시합니다.\n"
        "Vercel 서버리스는 번들 250 MB 제한이라 PyTorch 를 싣지 못합니다. "
        "ONNX Runtime 을 직접 호출하는 경로가 현실적입니다 — 순전파가 25.9 ms 이므로 항만 한 곳이 1 초 아래로 내려갑니다.\n"
        "탐지 결과는 회전 박스가 아니라 중심점과 길이·방향으로 표시합니다. 10 m 에서 실제로 측정되는 것이 그 셋이기 때문입니다.\n"
        "한 곳이 자리를 잡으면 같은 조건을 만족하는 항만부터 차례로 넓힙니다. "
        "탁수 해역은 NIR(B08) 을 4번째 채널로 넣은 뒤에 엽니다 — 탁한 물도 근적외에서는 어둡기 때문입니다.",
        SANS, 9.5, BODY, line=1.52)


def build(out, fig):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    slide1(prs, fig)
    slide2(prs, fig)
    slide3(prs, fig)
    prs.save(out)
    print("저장:", out)


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--outdir", default="outputs")
    ap.add_argument("--out", default="3차업무_정승원.pptx")
    a = ap.parse_args()
    build(a.out, lambda n: os.path.join(a.outdir, n))
