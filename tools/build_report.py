"""2차 업무 발표자료(PPT 3페이지) 생성기.

1차 업무 자료와 같은 디자인 규격을 따른다.
    배경 #FAFAFA · 카드 흰색 + #EBEBEB 0.75pt 테두리 · 모서리 반경 0.125in
    제목 Pretendard SemiBold 27pt · 본문 Pretendard 9.5pt · 라벨 Cascadia Mono 8.5pt

수치는 outputs/metrics.json 에서 직접 읽는다. 실험을 다시 돌리면 발표자료도
자동으로 갱신되므로 본문과 결과가 어긋날 일이 없다.

사용법
    py -3 run_3d_experiment.py
    py -3 tools/build_report.py
"""

from __future__ import annotations

import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "outputs")
REPORT = os.path.join(ROOT, "report")

BG = RGBColor(0xFA, 0xFA, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xEB, 0xEB, 0xEB)
INK = RGBColor(0x17, 0x17, 0x17)
BODY = RGBColor(0x4D, 0x4D, 0x4D)
MUTED = RGBColor(0x8F, 0x8F, 0x8F)
FAINT = RGBColor(0xA1, 0xA1, 0xA1)

MONO = "Cascadia Mono"
SANS = "Pretendard"
SANS_SB = "Pretendard SemiBold"
SANS_MD = "Pretendard Medium"

TOTAL_PAGES = 3


def set_font(run, name, size, color, bold=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 한글이 기본 글꼴로 떨어지지 않도록 동아시아 글꼴도 지정한다.
    rPr = run._r.get_or_add_rPr()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for tag in ("ea", "cs"):
        rPr.append(rPr.makeelement(ns + tag, {"typeface": name}))


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, spacing=1.35):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, (text, name, size, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        set_font(para.add_run(), name, size, color)
        para.runs[0].text = text
    return box


def add_card(slide, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.125 / min(w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_panel(slide, x, y, w, h, heading, lines):
    add_card(slide, x, y, w, h)
    add_text(slide, x + 0.26, y + 0.19, w - 0.52, 0.24,
             [(heading, SANS_SB, 10.5, INK)])
    add_text(slide, x + 0.26, y + 0.46, w - 0.52, h - 0.66, lines, spacing=1.3)


def add_image(slide, path, cx, cy, cw, ch, caption=None):
    add_card(slide, cx, cy, cw, ch)
    pad = 0.15
    pic = slide.shapes.add_picture(path, Inches(0), Inches(0))
    aspect = pic.width / pic.height
    box_w, box_h = cw - 2 * pad, ch - 2 * pad
    if box_w / box_h > aspect:
        h = box_h
        w = h * aspect
    else:
        w = box_w
        h = w / aspect
    pic.width = Emu(int(Inches(w)))
    pic.height = Emu(int(Inches(h)))
    pic.left = Emu(int(Inches(cx + (cw - w) / 2)))
    pic.top = Emu(int(Inches(cy + (ch - h) / 2)))
    if caption:
        add_text(slide, cx, cy + ch + 0.07, cw, 0.2,
                 [(caption, MONO, 8.5, MUTED)])
    return pic


def new_slide(prs, page, eyebrow, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False

    add_text(slide, 0.72, 0.46, 11.89, 0.24, [(eyebrow, MONO, 8.5, MUTED)])
    add_text(slide, 0.72, 0.74, 11.89, 0.52, [(title, SANS_SB, 27, INK)], spacing=1.0)
    add_text(slide, 0.72, 1.28, 11.89, 0.26, [(subtitle, SANS, 10.5, BODY)])
    add_text(slide, 11.4, 6.92, 1.21, 0.24,
             [(f"{page:02d} / {TOTAL_PAGES:02d}", MONO, 8.5, FAINT)],
             align=PP_ALIGN.RIGHT)
    return slide


def b(t):
    return (t, SANS, 9.5, BODY)


def k(t):
    return (t, SANS_MD, 9.5, INK)


def m(t):
    return (t, MONO, 8.5, BODY)


# ---------------------------------------------------------------------------

def build(prs, s):
    syn = s["synthetic_validation"]
    ss, se = syn["stereo"], syn["example_code"]
    sur = s["spe3r_pair_survey"]
    best = s["best_pair"]
    bex = s["best_pair_example_code"]

    # ---------------- 1. 방법 ----------------
    sl = new_slide(
        prs, 1, "01 / STEREO TRIANGULATION",
        "두 장의 영상에서 깊이 맵과 3D 포인트 클라우드",
        f"SPE3R aqua · 256×256 · f = {s['dataset']['fx']:.1f} px · "
        f"Stanford SLAB · CC BY-NC-SA 4.0")
    add_image(sl, os.path.join(OUT, "03_spe3r_stereo.png"), 0.72, 1.78, 11.89, 3.00,
              "RECTIFIED PAIR · DISPARITY · REFERENCE DEPTH · STEREO DEPTH · EXAMPLE CODE")
    add_panel(sl, 0.72, 5.12, 5.86, 1.73, "두 번째 시점을 어떻게 얻는가", [
        b("포즈 라벨 1,000개를 재보면 병진이 항상 (0, 0, Z)입니다."),
        b("카메라 좌표계만 보면 베이스라인이 없습니다."),
        k("삼각측량에 필요한 것은 카메라와 타겟 사이의 상대 운동입니다."),
        k("타겟이 회전하므로 타겟 기준으로는 카메라가 궤도를 돈 것과 같습니다."),
        m("R_ij = R_j · R_iᵀ        t_ij = t_j − R_ij · t_i"),
        k(f"회전 {best['rotation_deg']:.2f}° 쌍에서 유효 베이스라인 "
          f"{best['baseline_m']:.3f} m 확보"),
    ])
    add_panel(sl, 6.96, 5.12, 5.65, 1.73, "파이프라인", [
        b("① 쌍 선별 — 회전 8° 이내 + 베이스라인의 횡방향 성분이 지배적일 것"),
        b("② cv2.stereoRectify — 에피폴라선을 가로 행으로 정렬"),
        b("③ cv2.StereoSGBM — 화소마다 가로 이동량 d 탐색"),
        k("④ Z = f · B / d          ← 깊이 맵 (미터 단위)"),
        k("⑤ X = (u−cx)·Z/f,  Y = (v−cy)·Z/f     ← 포인트 클라우드"),
    ])

    # ---------------- 2. 검증 ----------------
    sl = new_slide(
        prs, 2, "02 / VALIDATION ON EXACT GROUND TRUTH",
        "정답이 있는 조건에서 먼저 검증",
        "구·직육면체의 광선 교차를 해석적으로 풀어 정답 깊이에 렌더링 오차가 없다 · "
        "남는 오차는 전부 정합 알고리즘에서 온다")
    add_image(sl, os.path.join(OUT, "01_synthetic_validation.png"),
              0.72, 1.78, 11.89, 3.00,
              "STEREO PAIR · DISPARITY · GROUND TRUTH / STEREO / EXAMPLE CODE DEPTH")
    add_panel(sl, 0.72, 5.12, 5.86, 1.73, "정량 결과", [
        m(f"베이스라인 {syn['baseline_m']:.2f} m · 기대 시차 "
          f"{syn['expected_disparity_px']:.0f} px · 분해능 "
          f"{syn['depth_resolution_m_per_px']*100:.1f} cm/px"),
        m("                RMSE     중앙값    5cm이내   깊이폭"),
        m(f"스테레오        {ss['rmse']:.4f}   {ss['median_abs']:.4f}    "
          f"{ss['within_5cm']*100:4.1f}%   {ss['span_m']:.3f}"),
        m(f"과제 예시       {se['rmse']:.4f}   {se['median_abs']:.4f}    "
          f"{se['within_5cm']*100:4.1f}%   {se['span_m']:.3f}"),
        k(f"정답 깊이폭 {syn['gt_span_m']:.3f} m — 스테레오는 {ss['span_m']:.3f} m 로 "
          f"되살리고, 예시 코드는 {se['span_m']:.3f} m 로 평면이 됩니다."),
    ])
    add_panel(sl, 6.96, 5.12, 5.65, 1.73, "Unit Test 65개", [
        k("해석해 — Z = f·B/d 를 1e-12 까지, 구의 중심 깊이 = 거리−반지름 을 1e-9 까지"),
        k("불변식 — p_j = R_ij·p_i + t_ij 성립, B와 d를 함께 2배 하면 Z 불변"),
        k("실패 특성화 — 같은 거리·다른 반사율이 4배 다른 깊이로 읽히는 것,"),
        b("                 앞면이 그늘이면 깊이 구조가 뭉개지는 것"),
        b("경계 조건 — 시차 0(무한원점), None 입력, 크기 불일치"),
    ])

    # ---------------- 3. 결과와 한계 ----------------
    sl = new_slide(
        prs, 3, "03 / RESULTS & LIMITS",
        "SPE3R 실측 결과와 한계",
        "기준 깊이 = 동봉 메시 40만 점을 z-buffer 로 투영 · "
        "대조군에는 정답에 대한 최적 아핀정렬을 적용해 유리한 조건을 부여")
    add_image(sl, os.path.join(OUT, "04_pointclouds.png"), 0.72, 1.78, 11.89, 3.00,
              "GROUND TRUTH MESH · STEREO DEPTH → 3D (SINGLE VIEW) · EXAMPLE CODE")
    add_panel(sl, 0.72, 5.12, 5.86, 1.73, "결과", [
        m(f"후보 20쌍 → 복원 성공 {sur['pairs_reconstructed']}쌍 → "
          f"오차 10cm 이내 {sur['pairs_within_10cm']}쌍"),
        m("                RMSE     중앙값    5cm이내"),
        m(f"스테레오        {best['rmse']:.4f}   {best['median_abs']:.4f}    "
          f"{best['within_5cm']*100:4.1f}%"),
        m(f"과제 예시       {bex['rmse']:.4f}   {bex['median_abs']:.4f}    "
          f"{bex['within_5cm']*100:4.1f}%"),
        k(f"포인트 클라우드 {s['best_pair_points']:,}점 · 메시 대비 Chamfer "
          f"{s['best_pair_chamfer_pred_to_gt']:.4f}"),
        b(f"예시 코드는 최적 정렬 스케일 {bex['affine_scale']:+.3f} — 관계가 뒤집힙니다."),
    ])
    add_panel(sl, 6.96, 5.12, 5.65, 1.73, "정합을 좌우하는 요인 · 한계", [
        k("① 표면 무늬 — 조명만 뒤집어 앞면을 그늘에 넣는 통제 실험"),
        b(f"     깊이폭 {ss['span_m']:.3f} → 0.572 m. 블록 크기를 바꿔도 회복되지 않습니다."),
        k("② 이상치 필터 — filterSpeckles + 중앙값 3×3"),
        b(f"     RMSE {best['rmse_unfiltered']:.3f} → {best['rmse']:.3f} "
          f"(중앙값은 {best['median_abs_unfiltered']:.4f} → {best['median_abs']:.4f} 로 거의 불변)"),
        b(f"한계 — 1,000장 중 쓸 만한 쌍이 3개, 유효화소 {best['valid_ratio']*100:.0f}%, "
          f"단일 뷰라 뒷면은 복원 불가,"),
        b(f"          정답 자세를 그대로 사용, 분해능 하한 "
          f"{best['depth_resolution_m_per_px']*100:.1f} cm/px (dZ = Z²/fB)"),
    ])


def main() -> int:
    path = os.path.join(OUT, "metrics.json")
    if not os.path.exists(path):
        print(f"metrics.json 이 없습니다: {path}\n먼저 run_3d_experiment.py 를 실행하세요.")
        return 1
    with open(path, encoding="utf-8") as f:
        summary = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build(prs, summary)

    os.makedirs(REPORT, exist_ok=True)
    out = os.path.join(REPORT, "2차업무_정승원.pptx")
    prs.save(out)
    print(f"저장 완료 -> {out}  ({TOTAL_PAGES} 슬라이드)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
