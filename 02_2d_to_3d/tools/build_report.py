"""2차 업무 발표자료(PPT 4페이지) 생성기.

    달 지형 스테레오 실험의 결과를 4장으로 정리한다.

1차 업무 자료와 같은 디자인 규격을 따른다.
    배경 #FAFAFA · 카드 흰색 + #EBEBEB 0.75pt 테두리 · 모서리 반경 0.125in
    제목 Pretendard SemiBold 27pt · 본문 Pretendard 9.5pt · 라벨 Cascadia Mono 8.5pt

수치는 outputs/metrics.json 에서 직접 읽는다. 실험을 다시 돌리면 발표자료도
자동으로 갱신되므로 본문과 결과가 어긋날 일이 없다.

생성물(.pptx)은 저장소에 커밋하지 않는다. 필요할 때 이 스크립트로 만든다.

사용법
    py -3 run_3d_experiment.py
    py -3 tools/build_report.py [출력_폴더]
"""

from __future__ import annotations

import io
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "outputs")
REPORT = os.path.join(ROOT, "report")

BG = RGBColor(0xFA, 0xFA, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xEB, 0xEB, 0xEB)
INK = RGBColor(0x17, 0x17, 0x17)
BODY = RGBColor(0x4D, 0x4D, 0x4D)
MUTED = RGBColor(0x8F, 0x8F, 0x8F)
FAINT = RGBColor(0xA1, 0xA1, 0xA1)

MONO = "Cascadia Mono"
SANS = "Pretendard"
SANS_SB = "Pretendard SemiBold"
SANS_MD = "Pretendard Medium"

TOTAL_PAGES = 4


def test_count() -> int:
    """실행된 테스트 개수를 pytest 리포트에서 읽는다.

    발표자료에 숫자를 직접 적어 두면 테스트를 늘렸을 때 조용히 어긋난다.
    """
    import re
    path = os.path.join(OUT, "pytest_report.txt")
    if not os.path.exists(path):
        return 0
    text = io.open(path, encoding="utf-8", errors="replace").read()
    hit = re.search(r"(\d+) passed", text)
    return int(hit.group(1)) if hit else 0


def josa(n, pair):
    """숫자 뒤에 붙는 조사를 읽는 소리에 맞춘다 ("3을" / "5를").

    수치를 f-string 으로 채우다 보니 값이 바뀌면 조사가 틀어진다. 한 자리
    끝소리가 받침으로 끝나는 수(0 영, 1 일, 3 삼, 6 육, 7 칠, 8 팔)와 10(십)
    뒤에는 앞쪽 조사를, 나머지(2 이, 4 사, 5 오, 9 구) 뒤에는 뒤쪽을 쓴다.
    """
    with_batchim, without = pair.split("/")
    last = int(n) % 10
    closed = last in (0, 1, 3, 6, 7, 8) or (int(n) % 100 == 10)
    # '으로/로' 만 규칙이 다르다. ㄹ 받침(1 일, 7 칠, 8 팔) 뒤에는 '로' 를 쓴다.
    if with_batchim == "으로" and last in (1, 7, 8):
        closed = False
    return f"{n}{with_batchim if closed else without}"


def code_lines(name="test_stereo_baseline_matches_the_convergence_formula"):
    """테스트 파일에서 함수 하나를 원문 그대로 읽어 온다.

    과제 요청이 "Unit Test 코드 및 실행 결과 문서화" 이므로 슬라이드에 코드를
    직접 싣는다. 슬라이드에 손으로 옮겨 적으면 코드가 바뀌었을 때 조용히
    어긋나므로 파일에서 읽는다.
    """
    path = os.path.join(ROOT, "tests", "test_terrain.py")
    src = io.open(path, encoding="utf-8").read().splitlines()
    start = next(i for i, l in enumerate(src) if l.startswith(f"def {name}("))
    end = start + 1
    while end < len(src) and not src[end].startswith("def "):
        end += 1
    while end > start and not src[end - 1].strip():
        end -= 1
    return src[start:end]


def test_counts() -> dict:
    """pytest 실행 결과에서 파일별 통과 수를 센다.

    슬라이드에 "기하 44 · 스테레오 27" 처럼 적으려면 이 값이 필요하다.
    손으로 적어 두면 테스트를 늘렸을 때 조용히 어긋난다.
    """
    import re
    from collections import Counter

    path = os.path.join(OUT, "pytest_report.txt")
    if not os.path.exists(path):
        return {}
    text = io.open(path, encoding="utf-8", errors="replace").read()
    hits = re.findall(r"^(\S+?\.py)::\S+ PASSED", text, re.M)
    return dict(Counter(f.split("/")[-1] for f in hits))


def report_lines():
    """pytest 실행 결과에서 슬라이드에 들어갈 만큼만 뽑는다.

    전문은 outputs/pytest_report.txt 에 그대로 있다. 여기서는 파일별 통과 수와
    테스트 이름 몇 개만 옮긴다. 이름을 고르는 기준은 '무엇을 검증하는지가
    이름만 봐도 드러나는 것' 이다.
    """
    import re
    from collections import Counter

    path = os.path.join(OUT, "pytest_report.txt")
    if not os.path.exists(path):
        return ["outputs/pytest_report.txt 가 없습니다"]
    text = io.open(path, encoding="utf-8", errors="replace").read()

    hits = re.findall(r"^(\S+?\.py)::(\S+) PASSED", text, re.M)
    per_file = Counter(f.split("/")[-1] for f, _ in hits)
    names = {n.split("[")[0] for n in (x[1] for x in hits)}
    total = re.search(r"(\d+) passed", text)

    picked = [n for n in (
        "test_nadir_view_of_flat_ground_has_constant_depth_equal_to_altitude",
        "test_stereo_baseline_matches_the_convergence_formula",
        "test_rendered_depth_recovers_the_original_elevation",
        "test_rendered_pair_aligns_at_the_disparity_the_geometry_predicts",
        "test_full_pipeline_recovers_the_terrain_within_one_pixel",
    ) if n in names]

    out = [f"collected {len(hits)} items", ""]
    for f, n in sorted(per_file.items()):
        out.append(f"{f:<24s}{n:>3d} passed")
    out += ["", "PASSED (발췌)"]
    out += [f"  {n}" for n in picked]
    out += ["  ...", "", f"{total.group(1) if total else '?'} passed"]
    return out


def set_font(run, name, size, color, bold=False, tracking=None):
    """글꼴을 지정한다. tracking 은 자간 [pt], 음수면 좁힌다.

    디자인 규격(DESIGN-vercel.md)의 Don't 에 "large heading 의 음수 자간을
    풀지 말 것" 이 있다. 표에 따르면 48px 제목이 -2.4px, 32px 제목이 -1.28px 로
    둘 다 글자 크기의 -4% 다. 같은 비율을 그대로 적용한다.
    """
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if tracking:
        # python-pptx 에 자간 API 가 없어 rPr 의 spc 속성을 직접 쓴다.
        # 단위는 1/100 pt 다.
        run._r.get_or_add_rPr().set("spc", str(int(round(tracking * 100))))
    # 한글이 기본 글꼴로 떨어지지 않도록 동아시아 글꼴도 지정한다.
    rPr = run._r.get_or_add_rPr()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for tag in ("ea", "cs"):
        rPr.append(rPr.makeelement(ns + tag, {"typeface": name}))


def style_paragraph(para, name, size, color, tracking=None):
    """문단과 그 안의 런에 모두 글꼴을 지정한다.

    런에만 지정하면 텍스트가 빈 문단은 기본 크기(18pt)로 높이를 잡아
    줄 간격이 어긋난다. 표에서 빈 칸이 한 행씩 밀리는 원인이었다.
    """
    para.font.name = name
    para.font.size = Pt(size)
    para.font.color.rgb = color
    run = para.add_run()
    set_font(run, name, size, color, tracking=tracking)
    return run


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, spacing=1.35,
             tracking=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, (text, name, size, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        style_paragraph(para, name, size, color, tracking).text = text
    return box


def add_card(slide, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.125 / min(w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_matrix(slide, x, y, col_w, rows, row_h=0.235, aligns=None,
               spacing=1.45):
    """표를 그린다. 열마다 별도 텍스트 상자를 둔다.

    등폭 글꼴로 공백을 맞추는 방식은 쓸 수 없다. Cascadia Mono 에는 한글이 없어
    대체 글꼴로 그려지는데 그 글꼴은 고정폭이 아니라서 열이 어긋난다.
    열을 상자로 분리하면 글꼴과 무관하게 정렬이 맞는다.

    Parameters
    ----------
    col_w : 열 너비 목록 [inch]. 기본은 첫 열만 왼쪽, 나머지는 오른쪽 정렬.
    rows : (텍스트 튜플, 스타일) 목록. 스타일은 'head' | 'key' | 'body' | 'small'.
    aligns : 열별 정렬 목록. 설명이 들어가는 열은 오른쪽으로 밀면 안 읽힌다.
    """
    styles = {
        "head": (SANS, 9, MUTED),
        "key": (SANS_MD, 10, INK),
        "body": (SANS, 10, BODY),
        "small": (SANS, 9, BODY),
        "small_key": (SANS_MD, 9, INK),
    }
    for j, cw in enumerate(col_w):
        left = x + sum(col_w[:j])
        box = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(cw),
                                       Inches(row_h * len(rows)))
        tf = box.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, (cells, style) in enumerate(rows):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = (aligns[j] if aligns else
                              PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT)
            para.line_spacing = spacing
            name, size, color = styles[style]
            style_paragraph(para, name, size, color).text = cells[j]


def add_panel(slide, x, y, w, h, heading, lines):
    add_card(slide, x, y, w, h)
    add_text(slide, x + 0.26, y + 0.19, w - 0.52, 0.24,
             [(heading, SANS_SB, 10.5, INK)])
    add_text(slide, x + 0.26, y + 0.46, w - 0.52, h - 0.66, lines, spacing=1.3)


def add_image(slide, path, cx, cy, cw, ch, caption=None):
    add_card(slide, cx, cy, cw, ch)
    pad = 0.15
    pic = slide.shapes.add_picture(path, Inches(0), Inches(0))
    aspect = pic.width / pic.height
    box_w, box_h = cw - 2 * pad, ch - 2 * pad
    if box_w / box_h > aspect:
        h = box_h
        w = h * aspect
    else:
        w = box_w
        h = w / aspect
    pic.width = Emu(int(Inches(w)))
    pic.height = Emu(int(Inches(h)))
    pic.left = Emu(int(Inches(cx + (cw - w) / 2)))
    pic.top = Emu(int(Inches(cy + (ch - h) / 2)))
    if caption:
        add_text(slide, cx, cy + ch + 0.07, cw, 0.2,
                 [(caption, MONO, 8.5, MUTED)])
    return pic


TALK = []
VARS = []


def add_notes(slide, text):
    """할 말을 모아 둔다. 슬라이드에는 넣지 않는다.

    발표자 노트로 넣으면 파일을 열어 본 사람에게 그대로 보이고, 인쇄하거나
    다른 도구로 변환할 때 따라다닌다. 발표자가 볼 것은 손에 든 대본이면
    충분하므로 pptx 는 슬라이드만 담고, 같은 내용을 대본 txt 로 따로 낸다.
    """
    TALK.append(text.strip())


def new_slide(prs, page, eyebrow, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False

    add_text(slide, 0.72, 0.46, 11.89, 0.24, [(eyebrow, MONO, 8.5, MUTED)])
    # 규격의 display 트래킹: 글자 크기의 -4%
    add_text(slide, 0.72, 0.74, 11.89, 0.52, [(title, SANS_SB, 27, INK)],
             spacing=1.0, tracking=-27 * 0.04)
    add_text(slide, 0.72, 1.28, 11.89, 0.26, [(subtitle, SANS, 10.5, BODY)])
    add_text(slide, 11.4, 6.92, 1.21, 0.24,
             [(f"{page:02d} / {TOTAL_PAGES:02d}", MONO, 8.5, FAINT)],
             align=PP_ALIGN.RIGHT)
    return slide


def b(t):
    """본문."""
    return (t, SANS, 10, BODY)


def k(t):
    """강조 — 수치나 결론."""
    return (t, SANS_MD, 10, INK)


def m(t):
    """등폭 — 표나 수식."""
    return (t, MONO, 9.5, BODY)


def t(x):
    """발표용 문장. 듣는 사람이 읽을 크기라 본문보다 한 단계 크다."""
    return (x, SANS, 11, BODY)


def tk(x):
    """발표용 문장 중 결론."""
    return (x, SANS_MD, 11, INK)


def gap():
    return ("", SANS, 4, BODY)


# ---------------------------------------------------------------------------


def build(prs, s):
    """4 페이지. 과제 결과물 형식(업무.pdf p.16)에 맞춰 자리를 나눈다."""
    sc, best = s["scene"], s["best"]
    conv = s["convergence_sweep"]
    blocks = s["block_size_sweep"]
    tc = test_count()
    res = sc["altitude_m"] ** 2 / (sc["focal_px"] * sc["baseline_m"])
    chosen = next(c for c in conv
                  if abs(c["convergence_deg"] - sc["convergence_deg"]) < 1e-9)
    widest = conv[-1]

    VARS[:] = [
        ("H", f"촬영 고도. {sc['altitude_m']/1000:.1f} km."),
        ("B", f"베이스라인 = 두 촬영 지점 사이 거리. "
              f"{sc['baseline_m']/1000:.1f} km."),
        ("수렴각", f"두 시선이 벌어진 각. {sc['convergence_deg']:.0f}도."),
        ("f", f"초점거리. {sc['focal_px']:,.0f} 픽셀."),
        ("d", "시차 = 같은 지점이 두 사진에서 옆으로 밀린 픽셀 수."),
        ("Z", "깊이 = 그 화소가 보는 지면까지의 거리."),
        ("고도", "Z 를 지형 좌표로 되돌린 값. 정답은 레이저 고도계 측정치."),
    ]

    # ---------------- 1. 방법 ----------------
    sl = new_slide(
        prs, 1, "01 / 방법",
        "달 사진 두 장으로 지형의 높낮이를 복원했습니다",
        f"티코 크레이터 · LOLA 레이저 고도 모델 {sc['gsd_m']/1000:.3f} km/화소 · "
        f"{sc['grid'][0]}×{sc['grid'][1]}")
    add_image(sl, os.path.join(OUT, "05_method.png"), 0.72, 1.70, 11.89, 3.16)

    # 왼쪽 — 식만. 라벨과 식을 열로 나눠야 식의 시작 위치가 맞는다.
    add_card(sl, 0.72, 5.02, 6.10, 1.86)
    add_text(sl, 0.98, 5.21, 5.58, 0.24, [("Z 를 뽑기까지", SANS_SB, 10.5, INK)])
    steps = [
        ("① 상대 자세", "R = R2 · R1^T,   t = t2 - R · t1"),
        ("② 베이스라인", "B = 2 · H · tan(수렴각/2)"),
        ("③ 정렬·정합", "stereoRectify(R, t) -> StereoSGBM -> d"),
        ("④ 깊이", "Z = f · B / d"),
        ("⑤ 3D 좌표", "X = (u-cx)·Z/f,   Y = (v-cy)·Z/f"),
    ]
    add_text(sl, 0.98, 5.52, 1.30, 0.94,
             [(k, SANS, 9.5, MUTED) for k, _ in steps], spacing=1.28)
    add_text(sl, 2.32, 5.52, 4.24, 0.94,
             [(v, MONO, 9.5, BODY) for _, v in steps], spacing=1.28)
    add_text(sl, 0.98, 6.58, 5.58, 0.22, [
        (f"→ 이렇게 나온 Z 오차 중앙값 {best['median_abs']/1000:.3f} km "
         f"(시차 {best['median_abs']/res:.2f} px)", SANS_MD, 9.5, INK)])

    # 오른쪽 — 기호가 무슨 역할인지와 그 값.
    add_card(sl, 7.05, 5.02, 5.56, 1.86)
    add_text(sl, 7.31, 5.21, 5.04, 0.24,
             [("각 기호가 하는 일", SANS_SB, 10.5, INK)])
    add_matrix(sl, 7.31, 5.55, [0.68, 2.86, 1.50], [
        (("H", "카메라가 떠 있는 높이",
          f"{sc['altitude_m']/1000:.1f} km"), "small"),
        (("B", "두 촬영 지점 사이 거리",
          f"{sc['baseline_m']/1000:.1f} km"), "small"),
        (("수렴각", "두 시선이 벌어진 각",
          f"{sc['convergence_deg']:.0f} 도"), "small"),
        (("f", "초점거리", f"{sc['focal_px']:,.0f} px"), "small"),
        (("d", "두 사진에서 밀린 픽셀 수", "찾는 값"), "small_key"),
        (("Z", "그 화소가 보는 지면까지 거리", "d 로 계산"), "small_key"),
    ], row_h=0.19, spacing=1.36,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT])

    add_notes(sl, f"""
· 달에 착륙하려면 어디가 평평한지 알아야 합니다. 사진은 색만 있고 높낮이가 없습니다.
· 그래서 같은 지역을 각도를 바꿔 두 번 찍고, 같은 지점이 두 사진에서 얼마나 밀렸는지로 거리를 계산했습니다.
· 촬영 고도 {sc['altitude_m']/1000:.0f} km, 두 촬영 지점 사이 {sc['baseline_m']/1000:.1f} km, 시선이 벌어진 각 {sc['convergence_deg']:.0f}도입니다.
· 나머지는 교과서대로입니다. 두 사진을 평행하게 편 뒤 시차 d 를 찾고 Z = f·B/d 로 나눕니다.
· 여기서 시차 1픽셀이 {res/1000:.2f} km 입니다. 이보다 미세한 높낮이는 원리적으로 구분되지 않습니다.
· 결과는 Z 오차 중앙값 {best['median_abs']/1000:.3f} km 입니다. 시차로 {best['median_abs']/res:.2f} 픽셀이니 부화소 수준입니다.
""")

    # ---------------- 2. Unit Test ----------------
    sl = new_slide(
        prs, 2, "02 / Unit Test",
        "크기와 자료형만 확인하면, 식이 틀려도 통과합니다",
        "실제 테스트 코드와 실행 결과 · tests/test_terrain.py · "
        "outputs/pytest_report.txt")

    add_card(sl, 0.72, 1.70, 7.35, 3.30)
    add_text(sl, 0.98, 1.89, 6.83, 0.24, [
        ("테스트 코드 — tests/test_terrain.py", SANS_SB, 10.5, INK)])
    add_text(sl, 0.98, 2.24, 6.83, 0.22, [
        ("베이스라인이 B = 2·H·tan(수렴각/2) 를 정확히 만족해야 합니다",
         SANS, 10, MUTED)])
    add_text(sl, 0.98, 2.58, 6.83, 2.30,
             [(x, MONO, 9.5, BODY) for x in code_lines()], spacing=1.3)

    add_card(sl, 8.31, 1.70, 4.30, 3.30)
    add_text(sl, 8.57, 1.89, 3.78, 0.24, [
        ("실행 결과 — pytest", SANS_SB, 10.5, INK)])
    add_text(sl, 8.57, 2.24, 3.78, 0.22, [
        ("걸린 시간만 지우고 그대로 저장합니다", SANS, 10, MUTED)])
    add_text(sl, 8.57, 2.58, 3.78, 2.20,
             [(x, MONO, 8.5, BODY) for x in report_lines()], spacing=1.2)

    add_card(sl, 0.72, 5.18, 11.89, 1.64)
    add_text(sl, 0.98, 5.62, 11.37, 0.90, [
        ("크기와 자료형만 보는 테스트는 계산식이 틀려도 통과합니다.",
         SANS, 11.5, BODY),
        ("그래서 손으로 답을 낼 수 있는 조건을 만들어 숫자까지 대조했습니다.",
         SANS, 11.5, BODY),
        (f"그 결과 {tc}개가 전부 통과했고, 렌더링이 시차를 밀어내던 버그를 "
         f"찾아 고쳤습니다.", SANS_MD, 11.5, INK),
    ], spacing=1.55)
    add_notes(sl, f"""
· 과제 예시 테스트는 크기와 자료형만 봅니다. 그러면 계산식이 틀려도 통과합니다.
· 그래서 손으로 답을 낼 수 있는 조건을 만들었습니다. 평지를 수직으로 내려다보면 모든 화소의 깊이가 촬영 고도와 같아야 합니다. 광축에 수직인 평면이기 때문입니다.
· 왼쪽이 그 테스트 코드입니다. 오른쪽이 실행 결과, {tc}개 전부 통과입니다.
· 이 테스트들이 실제로 버그를 잡았습니다. 지형을 그릴 때 점을 화소에 던지면서 위치를 반올림했는데, 그 오차가 경사 방향으로 쏠렸습니다.
· 두 촬영 지점은 서로 반대로 기울어 있어 좌우가 반대로 쏠렸고, 시차가 2.6 픽셀 치우쳐 고도가 1.245 km 틀렸습니다.
· 화소마다 광선을 쏘는 방식으로 바꿔 0.1 픽셀로 줄였습니다. 고도 오차는 {best['median_abs']/1000:.3f} km 가 됐습니다.
""")

    # ---------------- 3. 2D -> 3D 변환 결과 ----------------
    fu = s["fusion"]
    fz = fu["fused"]
    one = next(p for p in fu["per_angle"]
               if abs(p["convergence_deg"] - sc["convergence_deg"]) < 1e-9)
    widest = fu["per_angle"][-1]
    sl = new_slide(
        prs, 3, "03 / 2D → 3D 변환 결과",
        f"레이저로 잰 고도와 중앙값 {fz['median_abs']/1000:.3f} km 안에서 일치합니다",
        f"수렴각이 다른 네 쌍을 합친 결과 · 복원한 점 {s['n_points']:,}개")
    add_image(sl, os.path.join(OUT, "04_result.png"), 0.72, 1.70, 11.89, 3.16)

    add_card(sl, 0.72, 5.02, 4.30, 1.86)
    add_text(sl, 0.98, 5.21, 3.78, 0.24, [("결과", SANS_SB, 10.5, INK)])
    add_matrix(sl, 0.98, 5.60, [2.34, 1.44], [
        (("고도 오차 (중앙값)", f"{fz['median_abs']/1000:.3f} km"), "key"),
        (("90 퍼센타일", f"{fz['p90_abs']/1000:.3f} km"), "body"),
        (("지형 기복 대비", f"{fz['median_abs']/sc['relief_m']*100:.2f}%"), "key"),
        (("값이 나온 셀", f"{fz['coverage']*100:.0f}%"), "body"),
    ])
    add_panel(sl, 5.26, 5.02, 7.35, 1.86,
              "한 쌍보다 여러 쌍이 정확하고 넓습니다", [
        tk(f"기복 {sc['relief_m']/1000:.3f} km 인 크레이터를 오차 {fz['median_abs']/1000:.3f} km 로 "
           f"복원했습니다 — 기복의 {fz['median_abs']/sc['relief_m']*100:.2f}% 입니다."),
        t(f"한 쌍만 쓰면 {one['median_abs']/1000:.3f} km 에 {one['coverage']*100:.0f}%, "
          f"각도를 벌린 한 쌍은 {widest['median_abs']/1000:.3f} km 에 "
          f"{widest['coverage']*100:.0f}% 입니다. 합치면 둘 다 얻습니다."),
        t(f"남은 {100-fz['coverage']*100:.0f}% 는 크레이터 안쪽 그늘입니다. "
          f"두 사진에 같은 무늬가 보이지 않는 곳입니다."),
    ])
    add_notes(sl, f"""
· 왼쪽이 정답 고도, 가운데가 사진에서 복원한 고도, 오른쪽이 3D 점구름 {s['n_points']:,}개입니다.
· 크레이터 바닥과 테두리, 가운데 봉우리가 그대로 살아 있습니다.
· 숫자로는 고도 오차 중앙값 {fz['median_abs']/1000:.3f} km 입니다. 기복이 {sc['relief_m']/1000:.3f} km 인 지형이니 그 {fz['median_abs']/sc['relief_m']*100:.2f}퍼센트입니다.
· 정답은 레이저 고도계로 실제로 잰 값입니다. 제가 만든 근사가 아닙니다.
· 이건 수렴각이 다른 네 쌍을 합친 결과입니다. 각도가 다르면 정합에 성공하는 곳도 다릅니다.
· 한 쌍만 쓰면 오차 {one['median_abs']/1000:.3f} km 에 {one['coverage']*100:.0f}퍼센트, 각도를 크게 벌린 한 쌍은 {widest['median_abs']/1000:.3f} km 로 정확하지만 {widest['coverage']*100:.0f}퍼센트밖에 못 덮습니다.
· 합치니 오차도 줄고 덮는 범위도 넓어졌습니다. 정밀한 각도가 무거워지도록 분해능 제곱의 역수로 가중했습니다.
· 남은 {100-fz['coverage']*100:.0f}퍼센트는 크레이터 안쪽 그늘입니다.
""")

    # ---------------- 4. 개선점 ----------------
    # 고친 것은 3장의 수치에 이미 들어가 있다. 여기는 남은 것만 번호로 둔다.
    sl = new_slide(
        prs, 4, "04 / 개선점",
        "만들면서 부족하다고 느낀 것들",
        "앞줄이 무엇이 부족한지, 뒷줄이 왜 그것이 문제인지입니다")

    items = [
        ("1", "렌더링한 사진으로만 검증했다",
         "두 장 모두 고도 모델에서 제가 그린 것이라 그림자와 표면 반사가 "
         "실제 달 영상과 다릅니다.",
         "이대로면 이 숫자가 검증한 것은 촬영 기하와 정합이지 촬영 자체가 "
         "아닙니다."),
        ("2", "믿을 수 있는 곳을 정답 없이 가려내지 못한다",
         f"값이 나온 {fz['coverage']*100:.0f}% 안에서 어디가 맞았는지는 "
         f"정답 고도와 비교해야만 알 수 있습니다.",
         "실제 운용에는 정답이 없습니다. 각도 사이 편차를 대용으로 썼지만 "
         "그게 맞는지는 확인하지 못했습니다."),
        ("3", "파라미터를 지형마다 손으로 골라야 한다",
         "수렴각·블록 크기·대비 하한은 티코에서 훑어 고른 값입니다. "
         "코페르니쿠스에 그대로 쓰니 오차가 늘었습니다.",
         "지형이 바뀔 때마다 다시 고르는 것은 자동으로 도는 파이프라인이 "
         "아닙니다."),
        ("4", "시차가 정수 부근에 몰리는 편향이 남아 있다",
         "가로 2배로 정합해 49 → 36% 까지 낮췄지만 고르게 퍼졌을 때의 "
         "20% 보다 아직 높습니다.",
         "정합이 만든 편향을 뒤에서 덮고 있을 뿐, 원인을 정합 단계에서 "
         "없애지는 못했습니다."),
    ]
    add_card(sl, 0.72, 1.70, 11.89, 5.12)
    y = 2.10
    for num, head, now, why in items:
        add_text(sl, 0.98, y, 0.40, 0.28, [(num, SANS_SB, 15, FAINT)])
        add_text(sl, 1.42, y, 10.93, 0.28, [(head, SANS_SB, 14, INK)],
                 tracking=-14 * 0.02)
        add_text(sl, 1.42, y + 0.40, 10.93, 0.46,
                 [(now, SANS, 11, BODY), (why, SANS_MD, 11, INK)], spacing=1.3)
        y += 1.22
    add_notes(sl, f"""
· 개선점 넷입니다. 앞줄이 지금 잰 값이고 뒷줄이 할 일입니다.
· 1번이 가장 큽니다. 지금 쓰는 두 장은 고도 모델에서 제가 렌더링한 것입니다. 기하는 실제 달이지만 그림자와 표면 반사 특성은 실제 영상과 다릅니다.
· 그래서 이 실험이 검증한 것은 촬영 기하와 정합이지 촬영 자체가 아닙니다. 실제 궤도 영상을 넣어야 이 결과가 실제에서도 성립하는지 말할 수 있습니다.
· 2번은 그늘입니다. 값이 안 나오는 {100-fz['coverage']*100:.0f}퍼센트가 대부분 크레이터 안쪽인데, 두 사진에 같은 무늬가 안 보이면 원리적으로 맞출 수가 없습니다.
· 태양 각도가 다른 시점을 더 모으면 줄어듭니다. 실제 궤도 영상은 같은 지역을 여러 조명 조건에서 찍으니 가능한 방향입니다.
· 3번은 시차가 정수 부근에 몰리는 문제입니다. 가로로 2배 늘려 정합해서 49퍼센트에서 36퍼센트까지 낮췄는데, 균일하면 20퍼센트여야 하니 아직 남아 있습니다.
· 4배로 늘려 봤더니 오히려 나빠졌습니다. 보간이 없는 정보를 만들어 내지는 못하기 때문입니다.
· 4번은 고도 모델 해상도입니다. 지금은 0.118 km 화소라 매끈해서 잔무늬가 적습니다. 더 촘촘한 모델을 쓰면 정합이 쉬워지는데, 전용 포맷과 도구가 필요합니다.
""")


def talk_script(width=46):
    """모아 둔 할 말을 대본 txt 로 조립한다."""
    import textwrap

    titles = ["01 / 방법", "02 / Unit Test", "03 / 2D → 3D 변환 결과",
              "04 / 개선점"]
    out = ["2차 업무 발표 대본 — 정승원",
           f"슬라이드 {TOTAL_PAGES}장 · 각 장에서 할 말", "=" * 50, ""]
    if VARS:
        out += ["[0] 나오는 기호", "-" * 50]
        pad = max(len(k) for k, _ in VARS)
        for key, desc in VARS:
            wrapped = textwrap.wrap(desc, width + 14 - pad - 2)
            out.append(f"{key:<{pad}}  {wrapped[0]}")
            out += [" " * (pad + 2) + x for x in wrapped[1:]]
        out += ["", "쓰는 식", "  R = R2 · R1^T,    t = t2 − R · t1",
                "  B = 2 · H · tan(수렴각 / 2)", "  Z = f · B / d",
                "  X = (u−cx)·Z/f,    Y = (v−cy)·Z/f", ""]
    for i, note in enumerate(TALK):
        out += [f"[{i + 1}] {titles[i]}", "-" * 50]
        for line in note.splitlines():
            line = line.strip()
            if line:
                out += textwrap.wrap(line, width, subsequent_indent="  ")
        out.append("")
    return "\n".join(out)


def main() -> int:
    path = os.path.join(OUT, "metrics.json")
    if not os.path.exists(path):
        print(f"metrics.json 이 없습니다: {path}\n먼저 run_3d_experiment.py 를 실행하세요.")
        return 1
    with open(path, encoding="utf-8") as f:
        summary = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build(prs, summary)

    dest = sys.argv[1] if len(sys.argv) > 1 else REPORT
    os.makedirs(dest, exist_ok=True)
    out = os.path.join(dest, "2차업무_정승원.pptx")
    prs.save(out)
    print(f"저장 완료 -> {out}  ({TOTAL_PAGES} 슬라이드)")

    script = os.path.join(dest, "2차업무_발표대본_정승원.txt")
    io.open(script, "w", encoding="utf-8-sig",
            newline="\r\n").write(talk_script())
    print(f"저장 완료 -> {script}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
