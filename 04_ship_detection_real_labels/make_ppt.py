"""4차 업무 PPT 생성 — 코멘토 4차 업무 / 정승원

    py make_ppt.py --out 4차업무_정승원.pptx

2·3차 업무 PPT 서식을 그대로 따릅니다.
  13.333 x 7.5 in · 배경 #FAFAFA · 카드 #FFFFFF + #EBEBEB 0.75pt 테두리
  제목 Pretendard SemiBold 27pt · 본문 Pretendard · 수치 Cascadia Mono

구성
----
  01 자료    항만을 좌표로 식별하고 GSD 를 함급으로 검증
  02 결과    미국 해군기지 5곳 · 회전상자 · F1 0.936
  03 확인    탐지 결과를 항만마다 한 장씩
  04 부록    10 m 로 내리면 초해상화가 도움이 되는가 — 재학습이 이깁니다
"""
import os
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

W, H = 13.333, 7.5
INK, BODY, MUTE, FAINT = "171717", "4D4D4D", "8F8F8F", "A1A1A1"
HAIR, CANVAS, CARD = "EBEBEB", "FAFAFA", "FFFFFF"
ACC, WARN, OK = "0070F3", "D4443C", "1F8A4C"
SANS, SEMI, MED, MONO = ("Pretendard", "Pretendard SemiBold",
                         "Pretendard Medium", "Cascadia Mono")
TOTAL = 4
HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "outputs")


def C(h):
    return RGBColor.from_string(h)


def txt(sl, l, t, w, h, s, font=SANS, size=10, color=BODY,
        align=PP_ALIGN.LEFT, line=1.25):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = C(color)
    return tb


def card(sl, l, t, w, h, radius=0.045):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = C(CARD)
    sh.line.color.rgb = C(HAIR)
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if sh.has_text_frame:
        sh.text_frame.text = ""
    return sh


def base(prs, eyebrow, title, sub, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C(CANVAS)
    bg.line.fill.background()
    bg.shadow.inherit = False
    txt(sl, .72, .46, 11.89, .24, eyebrow, MONO, 8.5, MUTE)
    txt(sl, .72, .74, 11.89, .52, title, SEMI, 27, INK)
    txt(sl, .72, 1.28, 11.89, .26, sub, SANS, 10.5, BODY)
    txt(sl, 11.40, 6.92, 1.21, .24, "%02d / %02d" % (page, TOTAL),
        MONO, 8.5, FAINT, PP_ALIGN.RIGHT)
    return sl


def rows(sl, l, t, cols, data, widths, size=9.0):
    x = l
    for c, wd in zip(cols, widths):
        if c:
            txt(sl, x, t, wd, .2, c, MED, 8.5, MUTE)
        x += wd
    y = t + .26
    for r in data:
        x = l
        for ci, (v, wd) in enumerate(zip(r, widths)):
            col = BODY
            if isinstance(v, tuple):
                v, col = v
            txt(sl, x, y, wd, .2, str(v), SANS if ci == 0 else MONO, size, col)
            x += wd
        y += .245
    return y


def strip(sl, l, t, w, stats, gap=None):
    gap = gap or w / len(stats)
    x = l
    for lab, val in stats:
        col = INK
        if isinstance(val, tuple):
            val, col = val
        txt(sl, x, t, gap - .1, .22, lab, SANS, 9.0, MUTE)
        txt(sl, x, t + .30, gap - .1, .34, val, SEMI, 17, col)
        x += gap


def pic(sl, name, l, t, w=None, h=None):
    p = os.path.join(FIG, name)
    if not os.path.exists(p):
        txt(sl, l, t, w or 4, .3, "(그림 없음: %s)" % name, MONO, 9, WARN)
        return
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    sl.shapes.add_picture(p, Inches(l), Inches(t), **kw)


# ------------------------------------------------------------------ 슬라이드
def slide1(prs):
    sl = base(prs, "01 / 자료", "항만을 좌표로 식별하고, GSD 를 함급으로 검증했습니다",
              "HRSC2016 — 공식 문서는 'six famous harbors' 라고만 합니다. XML 좌표를 군집화해 6곳을 찾고 지형으로 확인했습니다.", 1)
    card(sl, .72, 1.80, 7.35, 4.85)
    pic(sl, "fig2_port_spotcheck.png", .90, 1.98, w=7.0)
    txt(sl, .90, 6.20, 7.0, .4,
        "노란 상자 = 사람이 그린 회전상자 정답. 항공모함·구축함·강습상륙함이 부두에 접안한 미 해군기지 5곳.",
        SANS, 8.5, MUTE, line=1.4)

    card(sl, 8.30, 1.80, 4.31, 2.30)
    txt(sl, 8.50, 1.95, 3.9, .22, "경도 부호가 XML 에 없습니다", SEMI, 11, INK)
    txt(sl, 8.50, 2.25, 3.95, 1.8,
        "32.6623, 117.1211 을 동경으로 읽으면 중국 내륙(안후이성).\n"
        "서경으로 읽으면 샌디에이고 해군기지 2 km 안.\n\n"
        "25 km 군집 7개, 무르만스크 2개를 합치면 6개 — 문서의 'six' 와 일치.\n"
        "라벨 1,070장 중 1,059장(99 %)이 미국 항만.",
        SANS, 9.0, BODY, line=1.45)

    card(sl, 8.30, 4.30, 4.31, 2.35)
    txt(sl, 8.50, 4.45, 3.9, .22, "XML 의 GSD 1.07 m 는 명목값입니다", SEMI, 11, INK)
    txt(sl, 8.50, 4.75, 3.95, 1.9,
        "그대로 믿으면 배 길이 중앙 360 m, 최대 987 m (실존 최대 458 m).\n"
        "타일 레벨 18 · 위도로 계산하면 약 0.45 m.\n"
        "선박 2,964척 장변 p99 = 347 m ↔ Nimitz 급 333 m  (+4 %)\n"
        "노퍽 341×47 px = 154×21 m ↔ Arleigh Burke 급 155×20 m",
        SANS, 9.0, BODY, line=1.45)
    return sl


def slide2(prs):
    sl = base(prs, "02 / 결과", "미국 해군기지 5곳, 회전상자로 F1 0.936",
              "YOLO11m-OBB · 학습 429장 / test 451장 (공식 분할, 학습에 안 쓴 영상) · 67 epoch · 학습 시드 3개 평균 · conf 0.25 · IoU 0.5", 2)
    card(sl, .72, 1.80, 11.89, 1.15)
    strip(sl, .95, 1.98, 11.5, [
        ("precision", "0.916"), ("recall", "0.957"), ("F1", ("0.936", ACC)),
        ("AP50", "0.970"), ("AP50-95", "0.770"), ("시드 편차 (F1)", "±0.004")])

    card(sl, .72, 3.15, 5.6, 3.5)
    txt(sl, .92, 3.30, 5.2, .22, "항만별 AP50 — 그 항만의 test 영상만으로", SEMI, 11, INK)
    rows(sl, .92, 3.65, ["항만", "test 선박", "AP50", ""],
         [["샌디에이고, CA", "683", "0.972", ""],
          ["노퍽, VA", "308", "0.981", ""],
          ["메이포트, FL", "160", "0.952", ""],
          ["에버렛, WA", "52", "0.974", ("표본 부족", WARN)],
          ["뉴포트, RI", "23", "0.933", ("표본 부족", WARN)]],
         [2.0, 1.0, 1.0, 1.2])
    txt(sl, .92, 5.35, 5.2, 1.2,
        "항만 간 차이가 작고 순서는 시드마다 바뀝니다. 에버렛·뉴포트는 신뢰구간이\n"
        "항만 간 차이보다 커서 참고값으로만 둡니다.",
        SANS, 9.0, MUTE, line=1.45)

    card(sl, 6.55, 3.15, 6.06, 3.5)
    txt(sl, 6.75, 3.30, 5.7, .22, "왜 회전상자인가 · 어떻게 학습했나", SEMI, 11, INK)
    txt(sl, 6.75, 3.65, 5.7, 2.9,
        "군함은 부두에 비스듬히 댑니다. 축정렬 상자는 이웃 배와 부두를 크게 포함해\n"
        "길이·폭을 못 재고 IoU 도 흐려집니다. 정답이 회전상자라 모델도 회전상자를\n"
        "내고, 길이·폭·방향이 표로 나옵니다.\n\n"
        "학습 중 validation 을 끄고 고정 epoch 의 last.pt 를 씁니다 — best.pt 는\n"
        "비교군마다 다른 epoch 을 골라 '학습량 동일' 조건을 깹니다.\n"
        "epochs 는 GPU 예산(T4 9 h)에 맞춰 100 → 67, 모든 런 동일.\n\n"
        "항만 단위 재분할(학습 샌디에이고·노퍽 → 시험 에버렛·뉴포트)로 누수 영향을\n"
        "따로 확인 — 결론 불변.",
        SANS, 9.0, BODY, line=1.45)
    return sl


def slide3(prs):
    sl = base(prs, "03 / 확인", "항만마다 한 장씩 — 탐지(초록)와 정답(노랑)",
              "전부 학습에 쓰지 않은 test 영상. 웹앱에서 항만을 고르고 ◀ 이전 / 다음 ▶ 으로 돌아가며 같은 것을 볼 수 있습니다.", 3)
    card(sl, .72, 1.80, 11.89, 4.85)
    pic(sl, "fig1_ports_detections.png", .90, 2.55, w=11.5)
    txt(sl, .90, 4.45, 11.5, 1.9,
        "샌디에이고 GT 11 / 탐지 12 — 부두에 나란한 함정을 회전상자가 각각 감쌉니다.   "
        "뉴포트 GT 3 / 탐지 1 — 계류 중인 민간 화물선 옆의 소형선 둘을 놓쳤습니다.\n"
        "웹앱은 정답이 있는 test 영상에서 맞힘·오탐·놓침을 IoU 0.5 로 세어 보여주고, "
        "GeoTIFF 를 올리면 회전상자를 GeoJSON 폴리곤으로 내보냅니다.",
        SANS, 9.0, MUTE, line=1.5)
    return sl


def slide4(prs):
    sl = base(prs, "04 / 부록", "10 m 로 내리면 초해상화가 도움이 되는가 — 재학습이 이깁니다",
              "같은 영상을 Sentinel-2 급 10 m 로 열화(PSF → 면적평균 1/22 → 잡음 → JPEG) · Native · Bicubic ×4 · Real-ESRGAN ×4 · 같은 조건 · 시드 3 · paired bootstrap", 4)
    card(sl, .72, 1.80, 6.1, 4.85)
    txt(sl, .92, 1.95, 5.8, .22, "F1 · 고정 conf 0.25 · test 1,228척", SEMI, 11, INK)
    rows(sl, .92, 2.30, ["입력", "10 m 재학습", "HR 학습 그대로"],
         [["HR 0.45 m (상한선)", ("0.936", ACC), "—"],
          ["Native LR", "0.683", ("0.052", WARN)],
          ["Bicubic ×4", ("0.690", OK), ("0.045", WARN)],
          ["Real-ESRGAN ×4", "0.678", ("0.260", OK)]],
         [2.3, 1.7, 1.7])
    txt(sl, .92, 3.75, 5.8, 2.8,
        "재학습 · Real-ESRGAN vs Native   AP50 −0.009 [−0.016, −0.003]  유의하게 나쁨\n"
        "HR 그대로 · Real-ESRGAN vs Native  AP50 +0.180 [+0.169, +0.201]  유의 · FP 3.7배\n\n"
        "① 탐지기를 10 m 로 다시 학습할 수 있으면 SR 은 무효 — 세 입력이 같습니다.\n"
        "② HR 탐지기를 그대로 써야 할 때만 SR 이 무너진 재현율(0.03)을 0.17 로 살립니다.\n"
        "③ 그래도 재학습(0.68)의 절반이 안 됩니다. 답은 'SR 을 붙여라' 가 아니라 '10 m 로 학습하라'.\n"
        "④ 10 m 에서 잃는 건 80 m 미만 소형선(재현율 0.15) — 어떤 SR 도 못 살리고, Real-ESRGAN 은 더 놓칩니다.",
        SANS, 9.0, BODY, line=1.5)
    card(sl, 7.05, 1.80, 5.56, 4.85)
    pic(sl, "fig4_pr_curves_10m.png", 7.25, 1.98, w=5.2)
    txt(sl, 7.25, 4.30, 5.2, 2.2,
        "왼쪽(10 m 재학습): 세 곡선이 겹칩니다.\n"
        "오른쪽(HR 학습 그대로): Real-ESRGAN 만 살아 있고 Native·Bicubic 은 재현율 0.1 을 못 넘습니다.\n\n"
        "DiffBIR 은 T4 로 50 h+ 라 제외. Real-ESRGAN 은 이 열화(blur+noise+JPEG)와 같은 가정으로\n"
        "학습된 순전파 모델이라 골랐습니다(1,068장 83초).",
        SANS, 9.0, MUTE, line=1.5)
    return sl


def build(out):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    for f in (slide1, slide2, slide3, slide4):
        f(prs)
    prs.save(out)
    print("저장:", out)


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=os.path.join(HERE, "4차업무_정승원.pptx"))
    a = ap.parse_args()
    build(a.out)
