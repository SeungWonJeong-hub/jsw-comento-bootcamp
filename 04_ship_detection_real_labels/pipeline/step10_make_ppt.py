"""4차 업무 PPT 생성 — 코멘토 4차 업무 / 정승원

    py pipeline/step10_make_ppt.py --out 4차업무_정승원.pptx

2·3차 업무 PPT 서식을 그대로 따릅니다.
  13.333 x 7.5 in · 배경 #FAFAFA · 카드 #FFFFFF + #EBEBEB 0.75pt 테두리
  제목 Pretendard SemiBold 27pt · 본문 Pretendard · 수치 Cascadia Mono

구성 — 전문용어 없이 씁니다
--------------------------
  01 데이터셋   정답이 있는 위성사진을 어떻게 구했나
  02 방법       사진 받기 → 정답 정리 → 학습 → 채점 → 웹앱
  03 결과       10척 중 9척 넘게 찾음 · 항구별 · 사진 5장
"""
import os
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

W, H = 13.333, 7.5
INK, BODY, MUTE, FAINT = "171717", "4D4D4D", "8F8F8F", "A1A1A1"
HAIR, CANVAS, CARD = "EBEBEB", "FAFAFA", "FFFFFF"
ACC, WARN = "0070F3", "D4443C"
SANS, SEMI, MED, MONO = ("Pretendard", "Pretendard SemiBold",
                         "Pretendard Medium", "Cascadia Mono")
TOTAL = 3
HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))   # 04 폴더
FIG = os.path.join(HERE, "outputs")


def C(h):
    return RGBColor.from_string(h)


def txt(sl, l, t, w, h, s, font=SANS, size=10, color=BODY,
        align=PP_ALIGN.LEFT, line=1.25):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = C(color)
    return tb


def card(sl, l, t, w, h, radius=0.045):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = C(CARD)
    sh.line.color.rgb = C(HAIR)
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if sh.has_text_frame:
        sh.text_frame.text = ""
    return sh


def base(prs, eyebrow, title, sub, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C(CANVAS)
    bg.line.fill.background()
    bg.shadow.inherit = False
    txt(sl, .72, .46, 11.89, .24, eyebrow, MONO, 8.5, MUTE)
    txt(sl, .72, .74, 11.89, .52, title, SEMI, 27, INK)
    txt(sl, .72, 1.28, 11.89, .26, sub, SANS, 10.5, BODY)
    txt(sl, 11.40, 6.92, 1.21, .24, "%02d / %02d" % (page, TOTAL),
        MONO, 8.5, FAINT, PP_ALIGN.RIGHT)
    return sl


def rows(sl, l, t, cols, data, widths, size=9.0):
    x = l
    for c, wd in zip(cols, widths):
        if c:
            txt(sl, x, t, wd, .2, c, MED, 8.5, MUTE)
        x += wd
    y = t + .26
    for r in data:
        x = l
        for ci, (v, wd) in enumerate(zip(r, widths)):
            col = BODY
            if isinstance(v, tuple):
                v, col = v
            txt(sl, x, y, wd, .2, str(v), SANS if ci == 0 else MONO, size, col)
            x += wd
        y += .245
    return y


def strip(sl, l, t, w, stats, gap=None):
    gap = gap or w / len(stats)
    x = l
    for lab, val in stats:
        col = INK
        if isinstance(val, tuple):
            val, col = val
        txt(sl, x, t, gap - .1, .22, lab, SANS, 9.0, MUTE)
        txt(sl, x, t + .30, gap - .1, .34, val, SEMI, 17, col)
        x += gap


def pic(sl, name, l, t, w=None, h=None):
    p = os.path.join(FIG, name)
    if not os.path.exists(p):
        txt(sl, l, t, w or 4, .3, "(그림 없음: %s)" % name, MONO, 9, WARN)
        return
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    sl.shapes.add_picture(p, Inches(l), Inches(t), **kw)


def step(sl, x, w, n, title, body):
    """한 단계. 번호 · 제목 · 설명."""
    txt(sl, x, 2.20, w, .20, n, MONO, 9.0, ACC)
    txt(sl, x, 2.46, w, .24, title, SEMI, 11.5, INK)
    txt(sl, x, 2.80, w, 1.2, body, SANS, 9.0, MUTE, line=1.45)


# ------------------------------------------------------------------ 슬라이드
def slide1(prs):
    sl = base(prs, "01 / 데이터셋", "정답이 있는 위성사진을 구했습니다 — 미국 군항 5곳",
              "HRSC2016: 위성사진 1,059장, 배 2,964척마다 사람이 직접 상자를 그려 둔 공개 연구용 자료입니다.", 1)

    # 어디서 왔나 — 출처 사슬
    card(sl, .72, 1.80, 11.89, 1.95)
    txt(sl, .92, 1.95, 11.5, .22, "어디서 가져왔나", SEMI, 11, INK)
    w = 11.5 / 3
    for k, (n, t, b) in enumerate([
            ("원본 사진", "구글어스 (Google Earth)",
             "미국 해군기지 5곳의 위성사진을 화면에서 모은 것.\n한 칸이 땅에서 약 0.45 m."),
            ("정답 상자", "중국 서북공업대 연구진 (HRSC2016, 2017)",
             "사진마다 배 위치를 사람이 직접 상자로 그려\n연구용으로 공개. 논문: Liu et al., ICPRAM 2017."),
            ("내려받은 곳", "Kaggle 의 공개 사본  guofeng/hrsc2016",
             "원본 배포처 사본. 8 GB 압축 파일.\n사진 1,680장 + 정답 파일 1,681개.")]):
        x = .92 + k * w
        txt(sl, x, 2.30, w - .3, .2, n, MONO, 9.0, ACC)
        txt(sl, x, 2.55, w - .3, .24, t, SEMI, 11, INK)
        txt(sl, x, 2.88, w - .3, .8, b, SANS, 9.0, MUTE, line=1.45)

    card(sl, .72, 3.95, 3.85, 2.7)
    txt(sl, .92, 4.10, 3.5, .22, "왜 이 자료인가", SEMI, 11, INK)
    txt(sl, .92, 4.42, 3.5, 2.1,
        "정답이 없으면 '몇 개 맞혔는지' 를 잴 수 없습니다.\n\n"
        "그래서 사람이 정답을 그려 둔 공개 자료 중, 군항처럼 배가 많고\n"
        "크기가 다양한 곳을 골랐습니다.",
        SANS, 9.0, BODY, line=1.5)

    card(sl, 4.74, 3.95, 3.85, 2.7)
    txt(sl, 4.94, 4.10, 3.5, .22, "어느 항구인지 직접 알아냈습니다", SEMI, 11, INK)
    txt(sl, 4.94, 4.42, 3.5, 2.1,
        "자료에는 항구 이름이 없고 위도·경도만 있습니다.\n\n"
        "좌표를 지도에 찍어 보니 샌디에이고 · 노퍽 · 메이포트 ·\n"
        "에버렛 · 뉴포트 해군기지였고, 사진을 눈으로 봐서 확인했습니다.",
        SANS, 9.0, BODY, line=1.5)

    card(sl, 8.76, 3.95, 3.85, 2.7)
    txt(sl, 8.96, 4.10, 3.5, .22, "사진 한 칸이 땅에서 몇 m 인지 확인", SEMI, 11, INK)
    txt(sl, 8.96, 4.42, 3.5, 2.1,
        "자료에 적힌 값(1.07 m)대로면 배가 360 m — 말이 안 됩니다.\n\n"
        "위도와 지도 확대 단계로 다시 계산하니 0.45 m. 이걸로 재면\n"
        "가장 긴 배가 항공모함 길이(333 m)와 맞습니다.",
        SANS, 9.0, BODY, line=1.5)
    return sl


def slide2(prs):
    sl = base(prs, "02 / 방법", "사진 받기 → 정답 정리 → 학습 → 채점 → 웹앱",
              "다섯 단계 전부 코드로 돌아갑니다. 학습만 무료 GPU(Kaggle)를 썼고 나머지는 노트북에서 했습니다.", 2)
    card(sl, .72, 1.80, 11.89, 2.35)
    w = 11.89 / 5
    step(sl, .95, w - .2, "01", "사진 받기",
         "Kaggle 에서 내려받아 압축을 풉니다.\n사진 1,680장, 정답 파일 1,681개.")
    step(sl, .95 + w, w - .2, "02", "정답 정리",
         "정답 상자를 학습 프로그램(YOLO)이 읽는\n형식으로 바꾸고, 100장을 그려서 맞는지 확인.")
    step(sl, .95 + 2 * w, w - .2, "03", "학습",
         "사진 429장으로 YOLO 를 가르칩니다.\n같은 조건으로 세 번 반복해 운이 아닌지 확인.")
    step(sl, .95 + 3 * w, w - .2, "04", "채점",
         "학습에 안 쓴 사진 451장으로만 채점.\n항구마다 따로도 채점.")
    step(sl, .95 + 4 * w, w - .2, "05", "웹앱",
         "항구를 고르고 버튼을 누르면\n사진마다 찾은 배와 정답을 겹쳐 보여줍니다.")

    card(sl, .72, 4.35, 3.85, 2.3)
    txt(sl, .92, 4.50, 3.5, .22, "박스는 어떤 기준으로 치나", SEMI, 11, INK)
    txt(sl, .92, 4.80, 3.5, 1.8,
        "정답 상자: 배 한 척의 선체를 통째로 감싸는 가장 작은\n"
        "직사각형을, 배가 놓인 방향대로 기울여서 그림 (사람이).\n\n"
        "프로그램 상자: 같은 방식으로 그리되, 확신도 25 % 이상인\n"
        "것만 남기고 겹치는 상자는 하나로 합칩니다.",
        SANS, 9.0, BODY, line=1.5)
    card(sl, 4.74, 4.35, 3.85, 2.3)
    txt(sl, 4.94, 4.50, 3.5, .22, "맞혔다고 치는 기준", SEMI, 11, INK)
    txt(sl, 4.94, 4.80, 3.5, 1.8,
        "프로그램 상자와 정답 상자가 절반 이상 겹치면(IoU 0.5) 맞힘.\n"
        "정답 하나에 상자 하나만 인정 — 같은 배를 두 번 찾으면\n"
        "두 번째는 오탐.\n\n"
        "채점은 학습에 안 쓴 사진 451장으로만 합니다.",
        SANS, 9.0, BODY, line=1.5)
    card(sl, 8.76, 4.35, 3.85, 2.3)
    txt(sl, 8.96, 4.50, 3.5, .22, "왜 기울인 상자인가", SEMI, 11, INK)
    txt(sl, 8.96, 4.80, 3.5, 1.8,
        "군함은 부두에 비스듬히 대어 있습니다. 똑바른 네모로 감싸면\n"
        "옆 배와 부두까지 들어가서 길이를 못 잽니다.\n\n"
        "배 방향대로 기울이면 배만 딱 감싸고 길이·폭·방향까지 나옵니다.",
        SANS, 9.0, BODY, line=1.5)
    return sl


def slide3(prs):
    sl = base(prs, "03 / 결과", "10척 중 9척 넘게 찾습니다 — 항구 5곳 모두",
              "학습에 안 쓴 사진 451장(배 1,228척)으로 채점. 세 번 학습한 평균이고, 반복마다 차이는 ±0.004 뿐입니다.", 3)
    card(sl, .72, 1.80, 11.89, 1.15)
    strip(sl, .95, 1.98, 11.5, [
        ("찾은 것 중 진짜 배", "91.6 %"), ("실제 배 중 찾은 것", "95.7 %"),
        ("F1 (둘의 조화평균)", ("0.936", ACC)), ("AP50", "0.970"), ("AP50-95", "0.770")])

    card(sl, .72, 3.15, 4.9, 3.5)
    txt(sl, .92, 3.30, 4.6, .22, "항구별", SEMI, 11, INK)
    rows(sl, .92, 3.65, ["항구", "배", "정확도", "찾은 비율", "F1", "AP50"],
         [["샌디에이고", "683", "92 %", "95 %", "0.937", "0.966"],
          ["노퍽", "308", "95 %", "96 %", "0.955", "0.986"],
          ["메이포트", "160", "93 %", "95 %", "0.938", "0.940"],
          ["에버렛", "52", "84 %", "100 %", "0.912", "0.959"],
          ["뉴포트", "23", "95 %", "87 %", "0.909", "0.942"]],
         [1.2, .6, .75, .8, .7, .7])
    txt(sl, .92, 5.30, 4.6, 1.2,
        "정확도 = 찾은 것 중 진짜 배의 비율.  찾은 비율 = 실제 배 중 찾아낸 비율.\n"
        "오른쪽 사진: 노랑 = 사람이 그린 정답, 초록 = 프로그램이 친 상자.\n"
        "둘이 겹치면 '맞힘', 초록만 있으면 '오탐', 노랑만 있으면 '놓침'.",
        SANS, 9.0, MUTE, line=1.45)

    ports = [("fig1_san_diego.png", "샌디에이고 · 정답 11 · 맞힘 10 · 오탐 2 · 놓침 1"),
             ("fig1_norfolk.png", "노퍽 · 정답 7 · 맞힘 7 · 오탐 0 · 놓침 0"),
             ("fig1_mayport.png", "메이포트 · 정답 6 · 맞힘 5 · 오탐 0 · 놓침 1"),
             ("fig1_everett.png", "에버렛 · 정답 4 · 맞힘 4 · 오탐 1 · 놓침 0"),
             ("fig1_newport.png", "뉴포트 · 정답 3 · 맞힘 1 · 오탐 0 · 놓침 2")]
    x0, y0, cw, ch, gap = 5.85, 3.15, 2.18, 1.62, .12
    for k, (f, cap) in enumerate(ports):
        r, c = divmod(k, 3)
        x, y = x0 + c * (cw + gap), y0 + r * (ch + .22 + gap)
        card(sl, x, y, cw, ch + .22)
        pic(sl, f, x + .06, y + .06, w=cw - .12, h=ch - .12)
        txt(sl, x + .08, y + ch - .02, cw - .16, .2, cap, MED, 7.5, INK)
    return sl


def build(out):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    for f in (slide1, slide2, slide3):
        f(prs)
    prs.save(out)
    print("저장:", out)


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=os.path.join(HERE, "4차업무_정승원.pptx"))
    a = ap.parse_args()
    build(a.out)
